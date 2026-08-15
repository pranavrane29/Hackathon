from pathlib import Path
from io import BytesIO

import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

from langchain_text_splitters import RecursiveCharacterTextSplitter


# ---------------------------------------------------------
# TEXT SPLITTER
# ---------------------------------------------------------

text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=1000,
    chunk_overlap=200
)


# ---------------------------------------------------------
# HELPER: GET FILE NAME AND BYTES
# ---------------------------------------------------------

def get_file_data(file):
    """
    Supports both:
    - normal file paths
    - Streamlit UploadedFile objects

    Returns:
        file_name, file_bytes
    """

    # Normal path
    if isinstance(file, (str, Path)):

        path = Path(file)

        if not path.exists():
            raise FileNotFoundError(
                f"File not found: {path}"
            )

        return path.name, path.read_bytes()

    # Streamlit UploadedFile or similar object
    file_name = getattr(
        file,
        "name",
        "uploaded_file"
    )

    file_bytes = file.read()

    return file_name, file_bytes


# ---------------------------------------------------------
# PDF PROCESSING
# ---------------------------------------------------------

def process_pdf(file_name, file_bytes):
    """
    Extract text from PDF pages and create chunks.
    """

    reader = PdfReader(
        BytesIO(file_bytes)
    )

    chunks = []

    for page_number, page in enumerate(
        reader.pages,
        start=1
    ):

        text = page.extract_text()

        if not text or not text.strip():
            continue

        page_chunks = text_splitter.split_text(
            text
        )

        for chunk_number, chunk in enumerate(
            page_chunks,
            start=1
        ):

            chunks.append(
                {
                    "text": chunk,
                    "source": file_name,
                    "page": page_number,
                    "slide": None,
                    "row": None,
                    "chunk_id": (
                        f"{file_name}"
                        f"_page_{page_number}"
                        f"_chunk_{chunk_number}"
                    )
                }
            )

    return chunks


# ---------------------------------------------------------
# DOCX PROCESSING
# ---------------------------------------------------------

def process_docx(file_name, file_bytes):
    """
    Extract text from DOCX paragraphs.
    """

    document = Document(
        BytesIO(file_bytes)
    )

    paragraphs = []

    for paragraph in document.paragraphs:

        text = paragraph.text.strip()

        if text:
            paragraphs.append(text)

    full_text = "\n".join(paragraphs)

    if not full_text.strip():
        return []

    text_chunks = text_splitter.split_text(
        full_text
    )

    chunks = []

    for chunk_number, chunk in enumerate(
        text_chunks,
        start=1
    ):

        chunks.append(
            {
                "text": chunk,
                "source": file_name,
                "page": None,
                "slide": None,
                "row": None,
                "chunk_id": (
                    f"{file_name}"
                    f"_chunk_{chunk_number}"
                )
            }
        )

    return chunks


# ---------------------------------------------------------
# TXT PROCESSING
# ---------------------------------------------------------

def process_txt(file_name, file_bytes):
    """
    Extract text from a TXT file.
    """

    try:
        text = file_bytes.decode("utf-8")

    except UnicodeDecodeError:
        text = file_bytes.decode(
            "latin-1"
        )

    text = text.strip()

    if not text:
        return []

    text_chunks = text_splitter.split_text(
        text
    )

    chunks = []

    for chunk_number, chunk in enumerate(
        text_chunks,
        start=1
    ):

        chunks.append(
            {
                "text": chunk,
                "source": file_name,
                "page": None,
                "slide": None,
                "row": None,
                "chunk_id": (
                    f"{file_name}"
                    f"_chunk_{chunk_number}"
                )
            }
        )

    return chunks


# ---------------------------------------------------------
# CSV PROCESSING
# ---------------------------------------------------------

def process_csv(file_name, file_bytes):
    """
    Convert CSV rows into searchable text.
    """

    dataframe = pd.read_csv(
        BytesIO(file_bytes)
    )

    chunks = []

    for row_number, row in dataframe.iterrows():

        # Convert the complete row into readable text
        row_text = " | ".join(
            f"{column}: {row[column]}"
            for column in dataframe.columns
        )

        row_text = row_text.strip()

        if not row_text:
            continue

        row_chunks = text_splitter.split_text(
            row_text
        )

        for chunk_number, chunk in enumerate(
            row_chunks,
            start=1
        ):

            actual_row = row_number + 2

            chunks.append(
                {
                    "text": chunk,
                    "source": file_name,
                    "page": None,
                    "slide": None,
                    "row": actual_row,
                    "chunk_id": (
                        f"{file_name}"
                        f"_row_{actual_row}"
                        f"_chunk_{chunk_number}"
                    )
                }
            )

    return chunks


# ---------------------------------------------------------
# PPTX PROCESSING
# ---------------------------------------------------------

def process_pptx(file_name, file_bytes):
    """
    Extract text from PowerPoint slides.
    """

    presentation = Presentation(
        BytesIO(file_bytes)
    )

    chunks = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        slide_text = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        full_text = "\n".join(
            slide_text
        )

        if not full_text.strip():
            continue

        slide_chunks = text_splitter.split_text(
            full_text
        )

        for chunk_number, chunk in enumerate(
            slide_chunks,
            start=1
        ):

            chunks.append(
                {
                    "text": chunk,
                    "source": file_name,
                    "page": None,
                    "slide": slide_number,
                    "row": None,
                    "chunk_id": (
                        f"{file_name}"
                        f"_slide_{slide_number}"
                        f"_chunk_{chunk_number}"
                    )
                }
            )

    return chunks


# ---------------------------------------------------------
# PROCESS ONE DOCUMENT
# ---------------------------------------------------------

def process_single_document(file):
    """
    Detect file type and process one document.
    """

    file_name, file_bytes = get_file_data(
        file
    )

    extension = Path(
        file_name
    ).suffix.lower()

    if extension == ".pdf":

        return process_pdf(
            file_name,
            file_bytes
        )

    elif extension == ".docx":

        return process_docx(
            file_name,
            file_bytes
        )

    elif extension == ".txt":

        return process_txt(
            file_name,
            file_bytes
        )

    elif extension == ".csv":

        return process_csv(
            file_name,
            file_bytes
        )

    elif extension == ".pptx":

        return process_pptx(
            file_name,
            file_bytes
        )

    else:

        raise ValueError(
            f"Unsupported file format: {extension}"
        )


# ---------------------------------------------------------
# PROCESS MULTIPLE DOCUMENTS
# ---------------------------------------------------------

def process_documents(files):
    """
    Process multiple documents.

    Supports:
        PDF
        DOCX
        TXT
        CSV
        PPTX

    Returns:
        Combined list of chunks.
    """

    if not files:
        raise ValueError(
            "No documents were provided."
        )

    all_chunks = []

    for file in files:

        chunks = process_single_document(
            file
        )

        all_chunks.extend(chunks)

    return all_chunks